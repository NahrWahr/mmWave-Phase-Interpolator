#!/usr/bin/env python3
# v1.0 — UCI / HIE Lab PowerPoint theme generator
# Usage:  uv run presentation/build-theme.py
# Output: output/uci-theme.pptx
#
# Layouts (in order, visible in Google Slides / PowerPoint):
#   0  Title Slide   — centered title + subtitle + author, cream bg
#   1  Index         — gold title bar, section list, gold vertical accent
#   2  Content       — gold title bar + body text / bullets
#   3  Image + Text  — gold title bar, picture left (50%), text right (50%)
#   4  Image + Image — gold title bar, two picture placeholders (50/50)
#   5  Standout      — navy bg, large centered white text, no chrome

import subprocess
import uuid
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# ── Paths ─────────────────────────────────────────────────────────────────────
REPO     = Path(__file__).resolve().parent.parent
SVG_LOGO = REPO / 'presentation' / 'figures' / 'uci-logo.svg'
PNG_LOGO = REPO / 'presentation' / 'figures' / 'uci-logo.png'
OUTPUT   = REPO / 'output' / 'uci-theme.pptx'

# ── Palette (mirrors uci-theme.tex) ───────────────────────────────────────────
GOLD      = RGBColor(0xF0, 0xC0, 0x40)   # #F0C040
GOLD_PALE = RGBColor(0xFA, 0xF0, 0xCC)   # #FAF0CC
BLUE      = RGBColor(0x00, 0x64, 0xA4)   # #0064A4
NAVY      = RGBColor(0x1B, 0x3D, 0x6F)   # #1B3D6F
CREAM     = RGBColor(0xFF, 0xF8, 0xF0)   # #FFF8F0
AMBER     = RGBColor(0xC8, 0x83, 0x0A)   # #C8830A
SAND      = RGBColor(0xEE, 0xD8, 0xA0)   # #EED8A0
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide geometry (16:9) ──────────────────────────────────────────────────────
W        = Inches(13.333)
H        = Inches(7.5)
FOOTER_H = Inches(0.32)
TITLE_H  = Inches(0.72)
LOGO_H   = Inches(0.38)
BODY_TOP = TITLE_H
BODY_H   = H - TITLE_H - FOOTER_H
PAD      = Inches(0.25)

# ── XML namespaces ────────────────────────────────────────────────────────────
PNS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
ANS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

_c = [500]
def _sid() -> int:
    _c[0] += 1
    return _c[0]

def _hx(rgb: RGBColor) -> str:
    return f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'


# ── Theme XML ─────────────────────────────────────────────────────────────────

def patch_theme(prs: Presentation) -> None:
    """Replace color and font scheme in theme1.xml with UCI palette."""
    try:
        tp = prs.slide_master.part.part_related_by(RT.THEME)
    except KeyError:
        print('  warning: theme part not found — skipping color patch')
        return

    # python-pptx 1.x stores the theme as a raw blob; parse → modify → re-serialize
    e = etree.fromstring(tp._blob)

    # Color scheme
    clr = e.find(f'.//{{{ANS}}}clrScheme')
    clr.set('name', 'UCI HIE Lab')
    for tag, val in [
        ('dk1', _hx(NAVY)),      ('lt1', _hx(CREAM)),
        ('dk2', _hx(BLUE)),      ('lt2', _hx(GOLD_PALE)),
        ('accent1', _hx(GOLD)),  ('accent2', _hx(AMBER)),
        ('accent3', _hx(SAND)),  ('accent4', _hx(SAND)),
        ('accent5', _hx(BLUE)),  ('accent6', _hx(GOLD_PALE)),
        ('hlink', _hx(BLUE)),    ('folHlink', _hx(AMBER)),
    ]:
        el = clr.find(f'{{{ANS}}}{tag}')
        if el is not None:
            for c in list(el):
                el.remove(c)
            etree.SubElement(el, f'{{{ANS}}}srgbClr').set('val', val)

    # Font scheme — Fira Sans (matches Moloch/Metropolis default)
    fs = e.find(f'.//{{{ANS}}}fontScheme')
    if fs is not None:
        fs.set('name', 'UCI HIE Lab')
        for ft in ('majorFont', 'minorFont'):
            fe = fs.find(f'{{{ANS}}}{ft}')
            if fe is not None:
                lat = fe.find(f'{{{ANS}}}latin')
                if lat is None:
                    lat = etree.SubElement(fe, f'{{{ANS}}}latin')
                lat.set('typeface', 'Fira Sans')

    tp._blob = etree.tostring(e, xml_declaration=True, encoding='UTF-8', standalone=True)


# ── Low-level XML helpers ─────────────────────────────────────────────────────

def _st(obj) -> etree._Element:
    """Return the spTree element of a master/layout/slide object."""
    return obj.element.find(f'.//{{{PNS}}}spTree')


def _set_bg(obj, color: RGBColor) -> None:
    """Set a solid-fill background on a master, layout, or slide object."""
    cSld = obj.element.find(f'{{{PNS}}}cSld')
    bg = cSld.find(f'{{{PNS}}}bg')
    if bg is None:
        bg = etree.Element(f'{{{PNS}}}bg')
        cSld.insert(0, bg)
    else:
        for c in list(bg):
            bg.remove(c)
    bgPr = etree.SubElement(bg, f'{{{PNS}}}bgPr')
    sf   = etree.SubElement(bgPr, f'{{{ANS}}}solidFill')
    etree.SubElement(sf, f'{{{ANS}}}srgbClr').set('val', _hx(color))
    etree.SubElement(bgPr, f'{{{ANS}}}effectLst')


def _clear(layout) -> None:
    """Remove all drawing shapes from a layout's spTree."""
    st = _st(layout)
    for c in list(st):
        if etree.QName(c.tag).localname not in ('nvGrpSpPr', 'grpSpPr'):
            st.remove(c)


def _x(obj, xml: str) -> None:
    """Append a shape (XML string) to obj's spTree (top of z-stack)."""
    _st(obj).append(etree.fromstring(xml))


def _x0(obj, xml: str) -> None:
    """Prepend a shape at z=0 (behind all other shapes)."""
    _st(obj).insert(2, etree.fromstring(xml))


# ── XML shape builders ────────────────────────────────────────────────────────

def _rect(l, t, w, h, fill: str, name='r') -> str:
    return f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="{name}"/>
    <p:cNvSpPr/><p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(l)}" y="{int(t)}"/>
      <a:ext cx="{int(w)}" cy="{int(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
</p:sp>'''


def _txtbox(l, t, w, h, text: str, sz: int, bold: bool,
            color: str, algn='l', anchor='ctr', name='t') -> str:
    b = '1' if bold else '0'
    return f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="{name}"/>
    <p:cNvSpPr txBox="1"/><p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(l)}" y="{int(t)}"/>
      <a:ext cx="{int(w)}" cy="{int(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="{anchor}" wrap="square"/>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="{algn}"/>
      <a:r>
        <a:rPr lang="en-US" sz="{sz * 100}" b="{b}">
          <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
        </a:rPr>
        <a:t>{text}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def _sldnum() -> str:
    """Slide-number field placeholder for the master footer."""
    fid = '{' + str(uuid.uuid4()).upper() + '}'
    return f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="SlideNum"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="sldNum" sz="quarter" idx="11"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(Inches(0.12))}" y="{int(H - FOOTER_H)}"/>
      <a:ext cx="{int(Inches(1.5))}" cy="{int(FOOTER_H)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p>
      <a:fld id="{fid}" type="slidenum">
        <a:rPr lang="en-US" sz="900" b="1">
          <a:solidFill><a:srgbClr val="{_hx(NAVY)}"/></a:solidFill>
        </a:rPr>
        <a:t>&#x2039;#&#x203a;</a:t>
      </a:fld>
    </a:p>
  </p:txBody>
</p:sp>'''


def _title_ph(l, t, w, h, fill: str, center=False) -> str:
    ph_type = 'ctrTitle' if center else 'title'
    algn    = 'ctr'     if center else 'l'
    return f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="Title"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="{ph_type}"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(l)}" y="{int(t)}"/>
      <a:ext cx="{int(w)}" cy="{int(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr" lIns="342900"/>
    <a:lstStyle>
      <a:lvl1pPr algn="{algn}">
        <a:defRPr lang="en-US" sz="2400" b="1">
          <a:solidFill><a:srgbClr val="{_hx(NAVY)}"/></a:solidFill>
        </a:defRPr>
      </a:lvl1pPr>
    </a:lstStyle>
    <a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p>
  </p:txBody>
</p:sp>'''


def _subtitle_ph(l, t, w, h) -> str:
    return f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="Subtitle"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="subTitle" idx="1"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(l)}" y="{int(t)}"/>
      <a:ext cx="{int(w)}" cy="{int(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="t" lIns="342900"/>
    <a:lstStyle>
      <a:lvl1pPr algn="ctr">
        <a:defRPr lang="en-US" sz="1600">
          <a:solidFill><a:srgbClr val="{_hx(NAVY)}"/></a:solidFill>
        </a:defRPr>
      </a:lvl1pPr>
    </a:lstStyle>
    <a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p>
  </p:txBody>
</p:sp>'''


def _author_ph(l, t, w, h) -> str:
    """Date/author placeholder — repurposed to show author + institute."""
    return f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="Author"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="dt" idx="10"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(l)}" y="{int(t)}"/>
      <a:ext cx="{int(w)}" cy="{int(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr" lIns="342900"/>
    <a:lstStyle>
      <a:lvl1pPr algn="ctr">
        <a:defRPr lang="en-US" sz="1100" i="1">
          <a:solidFill><a:srgbClr val="{_hx(NAVY)}"/></a:solidFill>
        </a:defRPr>
      </a:lvl1pPr>
    </a:lstStyle>
    <a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p>
  </p:txBody>
</p:sp>'''


def _body_ph(l, t, w, h, idx=1, sz=1800) -> str:
    return f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="Body {idx}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph idx="{idx}"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(l)}" y="{int(t)}"/>
      <a:ext cx="{int(w)}" cy="{int(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="t" lIns="342900" rIns="342900" tIns="182880"/>
    <a:lstStyle>
      <a:lvl1pPr>
        <a:defRPr lang="en-US" sz="{sz}">
          <a:solidFill><a:srgbClr val="{_hx(NAVY)}"/></a:solidFill>
        </a:defRPr>
      </a:lvl1pPr>
    </a:lstStyle>
    <a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p>
  </p:txBody>
</p:sp>'''


def _pic_ph(l, t, w, h, idx=1) -> str:
    return f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="Picture {idx}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="pic" idx="{idx}"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(l)}" y="{int(t)}"/>
      <a:ext cx="{int(w)}" cy="{int(h)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{_hx(SAND)}"/></a:solidFill>
    <a:ln>
      <a:solidFill><a:srgbClr val="{_hx(NAVY)}"/></a:solidFill>
    </a:ln>
  </p:spPr>
</p:sp>'''


RNS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def _pic_shape(rId: str, l, t, cx, cy, sid: int, name: str) -> str:
    """Build a <p:pic> element referencing an image relationship."""
    return f'''<p:pic xmlns:p="{PNS}" xmlns:a="{ANS}" xmlns:r="{RNS}">
  <p:nvPicPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvPicPr/>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{rId}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(l)}" y="{int(t)}"/>
      <a:ext cx="{int(cx)}" cy="{int(cy)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''


def _add_logo(obj, logo_png: Path) -> None:
    """
    Add the UCI logo PNG to the OPC part and insert a <p:pic> into obj's spTree.
    Works for both SlideMasterPart and SlideLayoutPart.
    """
    if not logo_png.exists():
        return
    img_part, rId = obj.part.get_or_add_image_part(str(logo_png))
    px_w, px_h = img_part.image.size
    logo_h = LOGO_H
    logo_w = int(logo_h * px_w / px_h)
    left = int(W - logo_w - Inches(0.08))
    top  = int(Inches(0.07))
    _x(obj, _pic_shape(rId, left, top, logo_w, logo_h, _sid(), 'UCI Logo'))


# ── Master setup ──────────────────────────────────────────────────────────────

def setup_master(master, logo_png: Path) -> None:
    """Cream background, gold footer bar, slide number, HIE Lab label, UCI logo."""
    _set_bg(master, CREAM)

    _x(master, _rect(0, H - FOOTER_H, W, FOOTER_H, _hx(GOLD), 'Footer Bar'))
    _x(master, _txtbox(
        W - Inches(1.6), H - FOOTER_H, Inches(1.5), FOOTER_H,
        'HIE Lab', sz=9, bold=True, color=_hx(NAVY), algn='r', name='HIE Lab',
    ))
    _x(master, _sldnum())
    _add_logo(master, logo_png)


# ── Layout builders ────────────────────────────────────────────────────────────

def _add_title_bar(layout) -> None:
    """Gold-pale bar at top (z=0) plus title placeholder (z=1)."""
    _x0(layout, _rect(0, 0, W, TITLE_H, _hx(GOLD_PALE), 'Title Bar'))
    _x(layout, _title_ph(
        Inches(0.2), 0, W - Inches(1.5), TITLE_H, fill=_hx(GOLD_PALE),
    ))


def make_title_layout(layout) -> None:
    layout.name = 'Title Slide'
    _clear(layout)
    # Thin gold separator between title and subtitle
    _x(layout, _rect(
        Inches(0.5), H * 0.555, W - Inches(1.0), Inches(0.04),
        _hx(GOLD), 'Gold Rule',
    ))
    _x(layout, _title_ph(
        Inches(0.5), H * 0.195, W - Inches(1.0), H * 0.34,
        fill=_hx(CREAM), center=True,
    ))
    _x(layout, _subtitle_ph(
        Inches(0.5), H * 0.60, W - Inches(1.0), Inches(0.90),
    ))
    _x(layout, _author_ph(
        Inches(0.5), H * 0.75, W - Inches(1.0), Inches(0.38),
    ))


def make_index_layout(layout) -> None:
    """
    Section index: title bar + list of section names.
    Usage: type sections into the body placeholder. Bold + color
    the active section (navy, 22pt); leave the rest at 16pt gray.
    A gold vertical bar on the left anchors the list visually.
    """
    layout.name = 'Index'
    _clear(layout)
    _add_title_bar(layout)
    # Gold vertical accent bar
    _x(layout, _rect(
        Inches(0.38), BODY_TOP + Inches(0.20),
        Inches(0.06), BODY_H - Inches(0.40),
        _hx(GOLD), 'Section Accent',
    ))
    # Body: larger default font for readability as a section list
    _x(layout, _body_ph(
        Inches(0.62), BODY_TOP + PAD,
        W - Inches(0.92), BODY_H - PAD * 2,
        idx=1, sz=2200,
    ))


def make_content_layout(layout) -> None:
    layout.name = 'Content'
    _clear(layout)
    _add_title_bar(layout)
    _x(layout, _body_ph(
        Inches(0.30), BODY_TOP + PAD,
        W - Inches(0.60), BODY_H - PAD * 2,
        idx=1,
    ))


def make_image_text_layout(layout) -> None:
    layout.name = 'Image + Text'
    _clear(layout)
    _add_title_bar(layout)
    half = W / 2
    _x(layout, _pic_ph(
        PAD, BODY_TOP + PAD,
        half - PAD * 1.5, BODY_H - PAD * 2,
        idx=1,
    ))
    _x(layout, _body_ph(
        half + PAD * 0.5, BODY_TOP + PAD,
        half - PAD * 1.5, BODY_H - PAD * 2,
        idx=2,
    ))


def make_image_image_layout(layout) -> None:
    layout.name = 'Image + Image'
    _clear(layout)
    _add_title_bar(layout)
    half = W / 2
    gap  = Inches(0.10)
    _x(layout, _pic_ph(
        PAD, BODY_TOP + PAD,
        half - PAD - gap / 2, BODY_H - PAD * 2,
        idx=1,
    ))
    _x(layout, _pic_ph(
        half + gap / 2, BODY_TOP + PAD,
        half - PAD - gap / 2, BODY_H - PAD * 2,
        idx=2,
    ))


def make_standout_layout(layout) -> None:
    """Full navy slide for Q&A / section breaks — no master chrome."""
    layout.name = 'Standout'
    layout.element.set('showMasterSp', '0')
    _clear(layout)
    _set_bg(layout, NAVY)
    _x(layout, f'''<p:sp xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:nvSpPr>
    <p:cNvPr id="{_sid()}" name="Standout Title"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="ctrTitle"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{int(Inches(1.0))}" y="{int(Inches(2.0))}"/>
      <a:ext cx="{int(W - Inches(2.0))}" cy="{int(Inches(3.5))}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle>
      <a:lvl1pPr algn="ctr">
        <a:defRPr lang="en-US" sz="3600" b="1">
          <a:solidFill><a:srgbClr val="{_hx(WHITE)}"/></a:solidFill>
        </a:defRPr>
      </a:lvl1pPr>
    </a:lstStyle>
    <a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p>
  </p:txBody>
</p:sp>''')


# ── Example slides ────────────────────────────────────────────────────────────

def add_examples(prs: Presentation) -> None:
    """Add one demonstration slide per layout."""
    layouts = prs.slide_master.slide_layouts

    # 1 — Title
    s = prs.slides.add_slide(layouts[0])
    for ph in s.placeholders:
        t = ph.placeholder_format.type.name
        if t == 'CENTER_TITLE':
            ph.text = 'mmWave Phase Interpolators'
            ph.text_frame.paragraphs[0].font.size = Pt(36)
        elif t == 'SUBTITLE':
            ph.text = 'Architectures for Phase Control at V/E Band'
        elif t == 'DATE':
            ph.text = 'Rahul Narwar · HIE Lab, University of California, Irvine'

    # 2 — Index  (active section bold navy; inactive items dimmed)
    s = prs.slides.add_slide(layouts[1])
    for ph in s.placeholders:
        if ph.placeholder_format.type.name == 'TITLE':
            ph.text = 'Overview'
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            sections = [
                ('First Principles', True),
                ('Phased-Array Beam Steering', False),
                ('Wireline / SerDes', False),
                ('Proposed Architecture', False),
            ]
            for i, (name, active) in enumerate(sections):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = name
                p.font.size = Pt(22 if active else 17)
                p.font.bold = active
                p.font.color.rgb = NAVY if active else RGBColor(0x80, 0x90, 0xA8)
                p.space_before = Pt(10)

    # 3 — Content
    s = prs.slides.add_slide(layouts[2])
    for ph in s.placeholders:
        if ph.placeholder_format.type.name == 'TITLE':
            ph.text = 'Content Slide'
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = 'Key point'
            for item in (
                'Supporting detail one',
                'Supporting detail two',
                'Supporting detail three',
            ):
                p = tf.add_paragraph()
                p.text = item
                p.level = 1

    # 4 — Image + Text
    s = prs.slides.add_slide(layouts[3])
    for ph in s.placeholders:
        if ph.placeholder_format.type.name == 'TITLE':
            ph.text = 'Image + Text Layout'
        elif ph.placeholder_format.idx == 2:
            ph.text = 'Caption or explanation goes here on the right side.'

    # 5 — Image + Image
    s = prs.slides.add_slide(layouts[4])
    for ph in s.placeholders:
        if ph.placeholder_format.type.name == 'TITLE':
            ph.text = 'Image + Image Layout'

    # 6 — Standout
    s = prs.slides.add_slide(layouts[5])
    for ph in s.placeholders:
        ph.text = 'Questions?'


# ── Entry point ───────────────────────────────────────────────────────────────

def main() -> None:
    print('Converting UCI logo SVG → PNG')
    subprocess.run(
        ['rsvg-convert', '-w', '440', str(SVG_LOGO), '-o', str(PNG_LOGO)],
        check=True,
    )

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print('Patching theme (colors + Fira Sans fonts)')
    patch_theme(prs)

    print('Configuring slide master')
    setup_master(prs.slide_master, PNG_LOGO)

    layouts = prs.slide_master.slide_layouts
    print(f'Building 6 layouts ({len(layouts)} slots available)')
    make_title_layout(layouts[0])
    make_index_layout(layouts[1])
    make_content_layout(layouts[2])
    make_image_text_layout(layouts[3])
    make_image_image_layout(layouts[4])
    make_standout_layout(layouts[5])
    for i in range(6, len(layouts)):
        layouts[i].name = f'(unused {i})'

    print('Adding example slides')
    add_examples(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f'Saved → {OUTPUT}')


if __name__ == '__main__':
    main()
